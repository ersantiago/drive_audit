import os
import sys

import docx2txt
import openpyxl
import PyPDF2
import odf

from pptx import Presentation
from xlrd import open_workbook
from odf import teletype, text


def check_flag(readfile, list_conf):
    ndachk = {}
    ndaf = []
    for conf in list_conf:
        if conf in readfile:
            ndachk[conf] = 'Fail'
            ndaf.append(conf)
        else:
            ndachk[conf] = 'Pass'
    ndafj = ','.join(ndaf)
    return ndafj

def read_docx(filepath):
    readfile = docx2txt.process(filepath).lower()
    return readfile

def read_xlsx(filepath):
    loadfile = openpyxl.load_workbook(filepath)
    sheets = loadfile.sheetnames
    cellsv = []
    for sheet1 in sheets:
        activewks = loadfile[sheet1]
        alldat = list(activewks.values)
        for row in alldat:
            rowdat = ','.join([str(x) for x in row if x is not None])
            cellsv.append(rowdat)
    readfile = ','.join(cellsv).lower()
    return readfile

def read_pptx(filepath):
    loadpptx = Presentation(filepath)
    slidesv = []
    for slide in loadpptx.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slidesv.append((shape.text).strip())
    readfile = ','.join(slidesv).lower()
    return readfile

def read_pdf(file_path):
    loadpdf = PyPDF2.PdfReader(file_path)
    pdfv = []
    maxpages = len(loadpdf.pages)
    if maxpages >= 50:
        readfile = 'limit exceeded, pages > 100'
    else:
        for i in range(maxpages):
            cpage = ((loadpdf.pages[i]).extract_text()).replace('\n', '').strip()
            pdfv.append(cpage)
        readfile = ','.join(pdfv).lower()
    return readfile

def read_generic(file_path):
    # Supports: .csv, .txt, .ini, .log, .py
    readfile = open(file_path, 'r').read().replace('\n', '').strip()
    return readfile

def read_xls(file_path):
    wkb = open_workbook(file_path)
    xlsv = []
    for sheet in wkb.sheets():
        nrows = sheet.nrows
        for i in range(nrows):
            data = sheet.row_values(i)
            rowv = ','.join(str(v) for v in data)
            xlsv.append(rowv)
    readfile = ','.join(xlsv)
    return readfile

def read_odf(file_path):
    loadfile = odf.opendocument.load(file_path)
    textelts = loadfile.getElementsByType(text.P)
    odfv = []
    for telt in textelts:
        xtr = teletype.extractText(telt)
        odfv.append(xtr)
    readfile = ','.join(odfv).lower()
    return readfile

# supported formats: xlsx, pptx, docx, pdf, txt, log, xls